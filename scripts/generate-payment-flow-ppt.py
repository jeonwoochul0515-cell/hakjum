"""
토스페이먼츠 가맹점 심사용 결제경로 PPT 생성기

각 슬라이드는 캡처 가이드 + 캡처 영역 자리표시자(placeholder)로 구성.
사용자가 실제 스크린샷을 자리표시자에 드래그-드롭하면 완성.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "학점나비_결제경로.pptx"

BRAND_NAVY = RGBColor(0x0F, 0x23, 0x4E)
BRAND_BLUE = RGBColor(0x32, 0x82, 0xF6)
BRAND_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
BRAND_GRAY = RGBColor(0x64, 0x74, 0x8B)
BRAND_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs, title, subtitle):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_NAVY

    tf_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), prs.slide_width - Inches(1.6), Inches(2.5)
    )
    tf = tf_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)


def add_step_slide(prs, step_num, total, title, url, description, capture_notes):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    header.line.fill.background()
    header.fill.solid()
    header.fill.fore_color.rgb = BRAND_NAVY

    # STEP 배지
    step_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(2.5), Inches(0.5)
    )
    step_tf = step_box.text_frame
    step_tf.margin_left = Pt(0)
    p = step_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"STEP {step_num} / {total}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BRAND_AMBER

    # 단계 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), prs.slide_width - Inches(1.0), Inches(0.6)
    )
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = BRAND_NAVY

    # URL 박스
    if url:
        url_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(1.95),
            Inches(8),
            Inches(0.4),
        )
        url_box.line.fill.background()
        url_box.fill.solid()
        url_box.fill.fore_color.rgb = BRAND_LIGHT
        url_tf = url_box.text_frame
        url_tf.margin_left = Pt(10)
        url_tf.margin_top = Pt(2)
        p = url_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"  URL  {url}"
        run.font.size = Pt(11)
        run.font.name = "Consolas"
        run.font.color.rgb = BRAND_GRAY

    # 좌측: 캡처 placeholder
    cap_left = Inches(0.5)
    cap_top = Inches(2.55)
    cap_width = Inches(6.5)
    cap_height = Inches(4.4)
    cap = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cap_left, cap_top, cap_width, cap_height
    )
    cap.line.color.rgb = BRAND_BLUE
    cap.line.width = Pt(2)
    cap.line.dash_style = 7  # DASH
    cap.fill.solid()
    cap.fill.fore_color.rgb = BRAND_LIGHT
    cap_tf = cap.text_frame
    cap_tf.word_wrap = True
    p = cap_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\n\n[ 스크린샷 영역 ]\n\n"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BRAND_BLUE
    p2 = cap_tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "이 도형을 삭제하고\n실제 화면 캡처 이미지를 붙여넣으세요"
    run2.font.size = Pt(12)
    run2.font.color.rgb = BRAND_GRAY

    # 우측: 설명 + 체크리스트
    desc_left = Inches(7.2)
    desc_top = Inches(2.55)
    desc_width = prs.slide_width - desc_left - Inches(0.4)

    # 설명
    desc_box = slide.shapes.add_textbox(desc_left, desc_top, desc_width, Inches(1.5))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    p = desc_tf.paragraphs[0]
    run = p.add_run()
    run.text = "단계 설명"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = BRAND_GRAY

    p2 = desc_tf.add_paragraph()
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = description
    run2.font.size = Pt(11)
    run2.font.color.rgb = BRAND_NAVY

    # 캡처 체크리스트
    chk_box = slide.shapes.add_textbox(
        desc_left, Inches(4.3), desc_width, Inches(2.6)
    )
    chk_tf = chk_box.text_frame
    chk_tf.word_wrap = True
    p = chk_tf.paragraphs[0]
    run = p.add_run()
    run.text = "캡처 시 포함 항목"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = BRAND_GRAY

    for note in capture_notes:
        p = chk_tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"☐  {note}"
        run.font.size = Pt(10)
        run.font.color.rgb = BRAND_NAVY


def add_section_slide(prs, title, subtitle):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_BLUE

    tf_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), prs.slide_width - Inches(1.6), Inches(2.0)
    )
    tf = tf_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 표지
    add_title_slide(
        prs,
        "학점나비 결제경로 안내",
        "토스페이먼츠 가맹점 심사용 / 사업자: 법률사무소 청송law (102-78-00061)",
    )

    # 섹션: 결제 전 단계
    add_section_slide(
        prs, "1부. 결제 전 단계", "서비스 진입 → 회원가입 → 이용권 선택"
    )

    steps = [
        # (step, title, url, description, capture_notes)
        (
            "메인 홈페이지 진입",
            "https://hakjum.school",
            "사용자가 학점나비 메인 홈페이지에 접속하여 서비스 소개를 확인하는 단계입니다. "
            "고교학점제 AI 맞춤 과목 추천 서비스의 가치 제안과 기능 소개를 보여줍니다.",
            [
                "도메인이 https://hakjum.school 인지 확인",
                "메인 비주얼/서비스 슬로건",
                "헤더의 로그인 버튼 노출",
            ],
        ),
        (
            "회원가입 / 로그인",
            "https://hakjum.school/login",
            "결제 진행을 위해 사용자가 로그인합니다. 이메일/비밀번호 또는 Google 소셜 로그인을 "
            "지원합니다. 비회원 결제는 제공하지 않습니다.",
            [
                "로그인 폼 (이메일/비밀번호)",
                "Google 소셜 로그인 버튼",
                "회원가입 페이지 링크",
            ],
        ),
        (
            "이용권 페이지 진입",
            "https://hakjum.school/subscription",
            "로그인 후 이용권 페이지에서 3개 요금제(무료 / 과목 선택 리포트 4,900원 / 올인원 패키지 7,900원)와 "
            "포함 기능, 환불 정책을 한 화면에서 확인할 수 있습니다.",
            [
                "3개 플랜 가격 (0원 / 4,900원 / 7,900원)",
                "각 플랜별 포함 기능 리스트",
                "환불 정책 안내문",
                "이용약관·환불정책 동의 체크박스",
                "토스페이먼츠 안전결제 표기",
            ],
        ),
        (
            "약관 동의 및 상품 선택",
            "https://hakjum.school/subscription",
            "이용약관 및 환불정책 동의 체크 후 원하는 이용권 버튼(예: '리포트 받기' 또는 "
            "'올인원 시작하기')을 클릭하여 결제를 시작합니다. 동의 전에는 결제 버튼이 비활성화됩니다.",
            [
                "약관·환불정책 동의 체크된 상태",
                "결제 버튼 활성화 (예: 리포트 받기 / 올인원 시작하기)",
                "선택한 상품 가격이 명확히 노출",
            ],
        ),
    ]

    total = len(steps) + 5  # 결제 단계 + 결제창 + 카드/인증 + 결과 + 환불
    for i, (title, url, desc, notes) in enumerate(steps, start=1):
        add_step_slide(prs, i, total, title, url, desc, notes)

    # 섹션: 결제 진행 단계
    add_section_slide(
        prs, "2부. 결제 진행 단계", "토스페이먼츠 결제창 → 카드정보 입력 → 인증"
    )

    payment_steps = [
        (
            "토스페이먼츠 결제창 호출",
            "https://hakjum.school/subscription (토스 결제창 오버레이)",
            "결제 버튼 클릭 시 토스페이먼츠 SDK가 호출되어 표준 결제창이 오버레이로 노출됩니다. "
            "주문번호, 상품명, 결제금액이 자동으로 결제창에 전달됩니다.",
            [
                "토스페이먼츠 결제창 표시",
                "주문명 (예: 학점나비 과목 선택 리포트)",
                "결제금액 (4,900원 또는 7,900원)",
                "결제수단 선택 옵션 (카드)",
            ],
        ),
        (
            "카드 정보 입력",
            "토스페이먼츠 결제창",
            "사용자가 카드 정보(카드번호, 유효기간, CVC, 비밀번호 앞 2자리, 생년월일/사업자번호)를 "
            "입력합니다. 카드사 ARS/앱 인증 단계로 이어집니다.",
            [
                "카드번호 입력 필드",
                "유효기간 / CVC 입력",
                "비밀번호 / 생년월일 입력",
                "결제 진행 버튼",
            ],
        ),
        (
            "카드사 인증",
            "각 카드사 인증창",
            "카드사가 제공하는 인증창에서 ARS·SMS·앱카드·간편결제 등 사용자가 선택한 방식으로 "
            "본인 인증을 완료합니다.",
            [
                "카드사 인증 화면",
                "인증 수단 선택 (ARS / 앱 / 간편결제 등)",
                "인증 완료 메시지",
            ],
        ),
    ]

    for j, (title, url, desc, notes) in enumerate(payment_steps, start=len(steps) + 1):
        add_step_slide(prs, j, total, title, url, desc, notes)

    # 섹션: 결제 완료 후
    add_section_slide(
        prs, "3부. 결제 완료 후", "성공 화면 → 환불 안내"
    )

    after_steps = [
        (
            "결제 성공 화면",
            "https://hakjum.school/subscription/success",
            "결제 승인이 완료되면 성공 페이지로 이동합니다. 서버에서 토스페이먼츠 confirm API를 호출하여 "
            "결제를 확정하고, 사용자 계정에 이용권이 부여됩니다.",
            [
                "결제 성공 메시지",
                "주문번호 / 결제 금액 / 상품명",
                "이용 시작 안내",
                "보고서/홈으로 이동 버튼",
            ],
        ),
        (
            "환불 정책 안내",
            "https://hakjum.school/refund-policy",
            "환불정책 페이지를 통해 환불 가능 조건(리포트 미열람·구매 7일 이내)과 환불 불가 조건(리포트 PDF 다운로드 후)을 "
            "사용자가 언제든 확인할 수 있습니다. 환불 문의는 고객센터로 접수됩니다.",
            [
                "환불 가능 조건",
                "환불 불가 조건",
                "환불 신청 방법 / 고객센터 연락처",
            ],
        ),
    ]

    for k, (title, url, desc, notes) in enumerate(
        after_steps, start=len(steps) + len(payment_steps) + 1
    ):
        add_step_slide(prs, k, total, title, url, desc, notes)

    # 마무리
    add_title_slide(
        prs,
        "감사합니다",
        "문의: support@hakjumnabi.com / 1660-4452",
    )

    prs.save(str(OUT))
    print(f"생성 완료: {OUT}")
    print(f"슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
